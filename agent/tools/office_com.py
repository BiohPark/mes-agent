"""MS Office 편집 엔진 (COM 자동화 + 라이브러리 폴백)

요구사항: 기존 Office 문서를 **열고·편집·저장**하는 고품질 기능.
- 1순위: 설치된 Word/Excel을 COM(win32com)으로 구동 → Office 자체 엔진으로 완전 충실도
  (서식 보존, 수식, 수정추적 수락, 메모 작성, PDF 내보내기)
- 2순위(COM 불가 시): python-docx / openpyxl 로 자동 폴백

핵심 제약: COM은 STA(단일 스레드 아파트)다. 모든 COM 호출은 CoInitialize 한 **하나의 스레드**에서
이뤄져야 한다. 브라우저(browser.py)의 greenlet 단일 스레드 패턴을 그대로 재사용한다.

안전 가드(품질 핵심):
- 편집 전 자동 백업(.bak) — 원본 손실 방지
- 입력이 진짜 OOXML(zip)인지 검증 후 편집
- 작업 결과에 사용된 엔진(com/docx/openpyxl)과 백업 경로를 명시
"""

import os
import json
import zipfile
import shutil
import atexit
import functools
from datetime import datetime
from concurrent.futures import ThreadPoolExecutor, TimeoutError as FutureTimeoutError
from pathlib import Path

import psutil

from agent.core.timeouts import office_com_timeout, timeout_error_text

# ── pywin32 가용성 탐지 ───────────────────────────────────────
try:
    import pythoncom  # noqa: F401
    import win32com.client  # noqa: F401
    _HAS_PYWIN32 = True
except Exception:
    _HAS_PYWIN32 = False


# ── 전용 STA 단일 스레드 ──────────────────────────────────────
def _com_thread_init():
    """워커 스레드에서 COM을 초기화한다(STA). pywin32 없으면 무시."""
    if _HAS_PYWIN32:
        try:
            pythoncom.CoInitialize()
        except Exception:
            pass


_com_executor = ThreadPoolExecutor(
    max_workers=1, thread_name_prefix="office-com", initializer=_com_thread_init
)

# 우리가 띄운 Office 인스턴스 PID만 추적 — 타임아웃 시 사용자가 직접 연 Office는 건드리지 않는다.
_tracked_pids: set[int] = set()


def _pids_by_name(image: str) -> set:
    try:
        return {p.pid for p in psutil.process_iter(["name"])
                if (p.info.get("name") or "").lower() == image.lower()}
    except Exception:
        return set()


def _track_new_pid(image: str, before: set) -> None:
    """Dispatch 전/후 PID 차집합으로 새로 뜬 우리 인스턴스를 기록(베스트에포트)."""
    try:
        new = _pids_by_name(image) - before
        for pid in new:
            _tracked_pids.add(pid)
    except Exception:
        pass


def _recover_stuck_com() -> None:
    """COM 호출이 타임아웃되면: 우리 인스턴스 PID만 강제 종료 + 싱글턴/executor 재생성.
    멈춘 워커를 버리고 새 executor를 만들어 다음 호출이 정상 동작하게 한다."""
    global _com_executor, _word_app, _excel_app, _ppt_app
    for pid in list(_tracked_pids):
        try:
            psutil.Process(pid).kill()
        except Exception:
            pass
    _tracked_pids.clear()
    _word_app = _excel_app = _ppt_app = None
    try:
        _com_executor.shutdown(wait=False)
    except Exception:
        pass
    _com_executor = ThreadPoolExecutor(
        max_workers=1, thread_name_prefix="office-com", initializer=_com_thread_init
    )


def _on_com_thread(fn, tool_name: str = "office"):
    """핸들러를 전용 COM 스레드에서 실행하되, OFFICE_COM_TIMEOUT 내에 안 끝나면
    멈춘 인스턴스를 정리하고 구조화된 타임아웃 오류를 돌려준다(무한 행 방지)."""
    @functools.wraps(fn)
    def wrapper(*args, **kwargs):
        future = _com_executor.submit(fn, *args, **kwargs)
        timeout = office_com_timeout()
        try:
            return future.result(timeout=timeout)
        except FutureTimeoutError:
            _recover_stuck_com()
            return timeout_error_text(tool_name, timeout, progressed=False)
    return wrapper


# ── COM 앱 싱글턴 (전용 스레드에서만 접근) ────────────────────
_word_app = None
_excel_app = None
_ppt_app = None


def _get_word():
    global _word_app
    import win32com.client as win32
    if _word_app is None:
        _before = _pids_by_name("WINWORD.EXE")
        _word_app = win32.dynamic.Dispatch("Word.Application")
        _word_app.Visible = False
        try:
            _word_app.DisplayAlerts = False
        except Exception:
            pass
        _track_new_pid("WINWORD.EXE", _before)
    return _word_app


def _get_excel():
    global _excel_app
    import win32com.client as win32
    if _excel_app is None:
        _before = _pids_by_name("EXCEL.EXE")
        _excel_app = win32.dynamic.Dispatch("Excel.Application")
        _excel_app.Visible = False
        _excel_app.DisplayAlerts = False
        # 열기 시 행 유발 대화상자 억제: 링크 업데이트 질문·매크로 보안 프롬프트 차단
        for _attr, _val in (("AskToUpdateLinks", False), ("AutomationSecurity", 3)):
            try:
                setattr(_excel_app, _attr, _val)
            except Exception:
                pass
        _track_new_pid("EXCEL.EXE", _before)
    return _excel_app


def _get_ppt():
    global _ppt_app
    import win32com.client as win32
    if _ppt_app is None:
        _before = _pids_by_name("POWERPNT.EXE")
        _ppt_app = win32.dynamic.Dispatch("PowerPoint.Application")
        _track_new_pid("POWERPNT.EXE", _before)
        # PowerPoint는 일부 작업에서 비가시 상태를 거부하므로 최소화로 띄운다
        try:
            _ppt_app.WindowState = 2  # ppWindowMinimized
        except Exception:
            pass
    return _ppt_app


def _quit_apps():
    """프로세스 종료 시 좀비 WINWORD.EXE/EXCEL.EXE/POWERPNT.EXE 방지."""
    global _word_app, _excel_app, _ppt_app
    for app in (_word_app, _excel_app, _ppt_app):
        try:
            if app is not None:
                app.Quit()
        except Exception:
            pass
    _word_app = None
    _excel_app = None
    _ppt_app = None


def _close_apps_on_com_thread():
    if _HAS_PYWIN32:
        try:
            _com_executor.submit(_quit_apps).result(timeout=10)
        except Exception:
            pass


atexit.register(_close_apps_on_com_thread)


# ── 공통 안전 헬퍼 ────────────────────────────────────────────
def _abspath(path: str) -> str:
    return os.path.abspath(os.path.expanduser(path))


def _validate_ooxml(path: str) -> str | None:
    """진짜 OOXML(zip)인지 검증. 문제 있으면 오류 메시지, 정상이면 None."""
    if not os.path.exists(path):
        return f"파일이 존재하지 않습니다: {path}"
    if not zipfile.is_zipfile(path):
        return (f"올바른 Office 파일이 아닙니다(OOXML zip 아님): {path}. "
                "마크다운/텍스트를 .docx로 잘못 저장한 파일일 수 있습니다.")
    return None


def _backup(path: str) -> str | None:
    """편집 전 타임스탬프 백업본을 만든다. 경로 반환."""
    p = Path(path)
    if not p.exists():
        return None
    bak = p.with_name(f"{p.stem}.{datetime.now():%Y%m%d_%H%M%S}{p.suffix}.bak")
    shutil.copy2(p, bak)
    return str(bak)


def _libre_to_pdf(path: str, pdf_path: str = "") -> str | None:
    """LibreOffice 헤드리스 PDF 폴백(설치돼 있으면). 경로 반환, 없으면 None."""
    try:
        from agent.tools.office_libre import to_pdf
        return to_pdf(path, pdf_path)
    except Exception:
        return None


def _no_com_msg(op: str) -> str:
    return json.dumps({
        "error": f"이 작업({op})은 설치된 MS Office(COM)가 필요합니다. "
                 "이 PC에서 Word/Excel COM을 사용할 수 없습니다.",
        "engine": "none",
    }, ensure_ascii=False)


# ── Word 편집 ─────────────────────────────────────────────────

def word_edit_text(path: str, find: str, replace: str,
                   match_case: bool = False, whole_word: bool = False) -> str:
    """기존 Word 문서에서 텍스트를 찾아 바꾸고 저장합니다(서식 보존).
    COM 우선, 불가 시 python-docx 폴백. 편집 전 자동 백업."""
    path = _abspath(path)
    err = _validate_ooxml(path)
    if err:
        return json.dumps({"error": err}, ensure_ascii=False)
    backup = _backup(path)

    # 1) COM 경로
    if _HAS_PYWIN32:
        try:
            word = _get_word()
            doc = word.Documents.Open(path)
            try:
                count_before = doc.Content.Text.count(find)
                f = doc.Content.Find
                f.ClearFormatting()
                f.Replacement.ClearFormatting()
                # win32com 동적 디스패치는 Find.Execute의 키워드 인자(특히 Replace)를
                # 제대로 바인딩하지 못한다. 반드시 완전 위치 인자로 호출한다.
                # (FindText, MatchCase, MatchWholeWord, MatchWildcards, MatchSoundsLike,
                #  MatchAllWordForms, Forward, Wrap, Format, ReplaceWith, Replace)
                # Wrap=1(wdFindContinue), Replace=2(wdReplaceAll)
                f.Execute(find, match_case, whole_word, False, False, False,
                          True, 1, False, replace, 2)
                doc.Save()
            finally:
                doc.Close()
            return json.dumps({
                "path": path, "engine": "com", "backup": backup,
                "find": find[:50], "replace": replace[:50],
                "occurrences": count_before,
                "message": "Word 찾아바꾸기 완료(서식 보존)",
            }, ensure_ascii=False)
        except Exception as e:
            # COM 실패 → 폴백 시도
            com_err = str(e)
    else:
        com_err = "pywin32/COM 미사용"

    # 2) python-docx 폴백 (단락·표 셀 단위 치환)
    try:
        import docx
        doc = docx.Document(path)
        n = 0

        def _replace_in_paragraph(par):
            nonlocal n
            if find in par.text:
                # 단락 전체 텍스트를 첫 run에 합쳐 치환(런 분할로 인한 누락 방지)
                full = par.text.replace(find, replace)
                n += par.text.count(find)
                for r in par.runs:
                    r.text = ""
                if par.runs:
                    par.runs[0].text = full
                else:
                    par.add_run(full)

        for par in doc.paragraphs:
            _replace_in_paragraph(par)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for par in cell.paragraphs:
                        _replace_in_paragraph(par)
        doc.save(path)
        return json.dumps({
            "path": path, "engine": "docx", "backup": backup,
            "occurrences": n, "com_error": com_err,
            "message": "Word 찾아바꾸기 완료(python-docx 폴백 — 일부 서식은 단순화될 수 있음)",
        }, ensure_ascii=False)
    except Exception as e2:
        return json.dumps({"error": f"COM/라이브러리 모두 실패: {com_err} | {e2}",
                           "backup": backup}, ensure_ascii=False)


def word_insert_text(path: str, text: str, after_anchor: str = "") -> str:
    """기존 Word 문서에 텍스트를 삽입하고 저장합니다.
    after_anchor가 있으면 그 텍스트 뒤에, 없으면 문서 끝에 삽입합니다.
    RAG/URL로 조사한 내용을 기존 보고서에 채워 넣을 때 사용. 편집 전 자동 백업."""
    path = _abspath(path)
    err = _validate_ooxml(path)
    if err:
        return json.dumps({"error": err}, ensure_ascii=False)
    backup = _backup(path)

    # 1) COM 경로
    if _HAS_PYWIN32:
        try:
            word = _get_word()
            doc = word.Documents.Open(path)
            try:
                if after_anchor:
                    rng = doc.Content
                    f = rng.Find
                    f.ClearFormatting()
                    found = f.Execute(after_anchor, False, False, False, False, False, True, 1)
                    if not found:
                        doc.Close(SaveChanges=0)
                        return json.dumps({"error": f"앵커 텍스트를 찾지 못했습니다: {after_anchor[:50]}",
                                           "backup": backup}, ensure_ascii=False)
                    rng.InsertAfter("\n" + text)
                else:
                    rng = doc.Content
                    rng.Collapse(0)  # wdCollapseEnd
                    rng.InsertAfter("\n" + text)
                doc.Save()
            finally:
                doc.Close()
            return json.dumps({"path": path, "engine": "com", "backup": backup,
                               "inserted_chars": len(text),
                               "message": "Word 텍스트 삽입 완료"}, ensure_ascii=False)
        except Exception as e:
            com_err = str(e)
    else:
        com_err = "pywin32/COM 미사용"

    # 2) python-docx 폴백 (문서 끝에만 추가)
    try:
        import docx
        doc = docx.Document(path)
        doc.add_paragraph(text)
        doc.save(path)
        return json.dumps({"path": path, "engine": "docx", "backup": backup,
                           "inserted_chars": len(text), "com_error": com_err,
                           "note": "폴백은 문서 끝에만 삽입합니다(after_anchor 미지원)",
                           "message": "Word 텍스트 삽입 완료(python-docx 폴백)"}, ensure_ascii=False)
    except Exception as e2:
        return json.dumps({"error": f"COM/라이브러리 모두 실패: {com_err} | {e2}",
                           "backup": backup}, ensure_ascii=False)


def word_export_pdf(path: str, pdf_path: str = "") -> str:
    """Word 문서를 PDF로 내보냅니다. (설치된 Word 필요)"""
    path = _abspath(path)
    err = _validate_ooxml(path)
    if err:
        return json.dumps({"error": err}, ensure_ascii=False)
    if not pdf_path:
        pdf_path = str(Path(path).with_suffix(".pdf"))
    pdf_path = _abspath(pdf_path)
    if _HAS_PYWIN32:
        try:
            word = _get_word()
            doc = word.Documents.Open(path)
            try:
                doc.ExportAsFixedFormat(OutputFileName=pdf_path, ExportFormat=17)  # wdExportFormatPDF
            finally:
                doc.Close(SaveChanges=0)
            return json.dumps({"path": pdf_path, "engine": "com",
                               "message": "PDF 내보내기 완료"}, ensure_ascii=False)
        except Exception as e:
            com_err = str(e)
    else:
        com_err = "pywin32/COM 미사용"
    # COM 불가/실패 → LibreOffice 헤드리스 폴백
    out = _libre_to_pdf(path, pdf_path)
    if out:
        return json.dumps({"path": out, "engine": "libreoffice",
                           "message": "PDF 내보내기 완료(LibreOffice 폴백)"}, ensure_ascii=False)
    return json.dumps({"error": f"PDF 내보내기 실패: {com_err}. LibreOffice도 사용 불가.",
                       "engine": "none"}, ensure_ascii=False)


def word_accept_all_changes(path: str) -> str:
    """Word 문서의 모든 수정추적(Track Changes)을 수락하고 저장합니다. (설치된 Word 필요)"""
    path = _abspath(path)
    err = _validate_ooxml(path)
    if err:
        return json.dumps({"error": err}, ensure_ascii=False)
    if not _HAS_PYWIN32:
        return _no_com_msg("수정추적 수락")
    backup = _backup(path)
    try:
        word = _get_word()
        doc = word.Documents.Open(path)
        try:
            doc.AcceptAllRevisions()
            doc.Save()
        finally:
            doc.Close()
        return json.dumps({"path": path, "engine": "com", "backup": backup,
                           "message": "모든 수정추적 수락 완료"}, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"error": str(e), "backup": backup}, ensure_ascii=False)


def word_add_comment(path: str, anchor_text: str, comment: str) -> str:
    """Word 문서에서 특정 텍스트에 검토 메모(Comment)를 추가합니다. (설치된 Word 필요)"""
    path = _abspath(path)
    err = _validate_ooxml(path)
    if err:
        return json.dumps({"error": err}, ensure_ascii=False)
    if not _HAS_PYWIN32:
        return _no_com_msg("메모 작성")
    backup = _backup(path)
    try:
        word = _get_word()
        doc = word.Documents.Open(path)
        try:
            # Range를 변수로 잡고 그 Find를 실행하면, 일치 시 Range가 찾은 범위로 재정의된다.
            rng = doc.Content
            f = rng.Find
            f.ClearFormatting()
            found = f.Execute(anchor_text, False, False, False, False, False, True, 1)
            if not found:
                doc.Close(SaveChanges=0)
                return json.dumps({"error": f"앵커 텍스트를 찾지 못했습니다: {anchor_text[:50]}",
                                   "backup": backup}, ensure_ascii=False)
            doc.Comments.Add(Range=rng, Text=comment)
            doc.Save()
        finally:
            doc.Close()
        return json.dumps({"path": path, "engine": "com", "backup": backup,
                           "message": "검토 메모 추가 완료"}, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"error": str(e), "backup": backup}, ensure_ascii=False)


# ── Excel 편집 ────────────────────────────────────────────────

def excel_set_cells(path: str, cells: dict, sheet: str = "") -> str:
    """기존 Excel 워크북의 특정 셀들에 값/수식을 설정하고 저장합니다.
    cells 예: {"B2": "=A1+A2", "C3": "합계"}. COM 우선, 불가 시 openpyxl 폴백."""
    path = _abspath(path)
    err = _validate_ooxml(path)
    if err:
        return json.dumps({"error": err}, ensure_ascii=False)
    if not isinstance(cells, dict) or not cells:
        return json.dumps({"error": "cells는 비어있지 않은 {셀주소: 값} 객체여야 합니다."},
                          ensure_ascii=False)
    backup = _backup(path)
    has_formula = any(isinstance(v, str) and v.startswith("=") for v in cells.values())

    # 1) COM 경로
    if _HAS_PYWIN32 and not has_formula:
        try:
            excel = _get_excel()
            # Notify=False: 파일이 잠겨 있으면 "사용 중" 대화상자로 무한 대기하지 않고 즉시 예외
            wb = excel.Workbooks.Open(
                path, UpdateLinks=0, IgnoreReadOnlyRecommended=True, Notify=False, AddToMru=False
            )
            try:
                ws = wb.Worksheets(sheet) if sheet else wb.ActiveSheet
                for addr, val in cells.items():
                    ws.Range(addr).Value = val
                wb.Save()
            finally:
                wb.Close(SaveChanges=0)
            return json.dumps({"path": path, "engine": "com", "backup": backup,
                               "cells_set": len(cells),
                               "message": "Excel 셀 편집 완료(수식·서식 보존)"}, ensure_ascii=False)
        except Exception as e:
            com_err = str(e)
    else:
        com_err = "수식 입력은 COM 계산 대기 없이 openpyxl 폴백" if has_formula else "pywin32/COM 미사용"

    # 2) openpyxl 폴백
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path)
        ws = wb[sheet] if sheet else wb.active
        for addr, val in cells.items():
            ws[addr] = val
        wb.save(path)
        return json.dumps({"path": path, "engine": "openpyxl", "backup": backup,
                           "cells_set": len(cells), "com_error": com_err,
                           "message": "Excel 셀 편집 완료(openpyxl 폴백 — 차트 등 일부 객체는 유실될 수 있음)"},
                          ensure_ascii=False)
    except Exception as e2:
        return json.dumps({"error": f"COM/라이브러리 모두 실패: {com_err} | {e2}",
                           "backup": backup}, ensure_ascii=False)


def excel_get_range(path: str, cell_range: str, sheet: str = "") -> str:
    """Excel 워크북의 특정 범위 값을 읽습니다. 예: 'A1:C10'. COM 우선, openpyxl 폴백."""
    path = _abspath(path)
    err = _validate_ooxml(path)
    if err:
        return json.dumps({"error": err}, ensure_ascii=False)

    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, data_only=True)
        ws = wb[sheet] if sheet else wb.active
        rows = [[c.value for c in row] for row in ws[cell_range]]
        if any(v is None for r in rows for v in r):
            wbf = openpyxl.load_workbook(path, data_only=False)
            wsf = wbf[sheet] if sheet else wbf.active
            frows = [[c.value for c in row] for row in wsf[cell_range]]
            rows = [
                [v if v is not None else fv for v, fv in zip(r, fr)]
                for r, fr in zip(rows, frows)
            ]
        return json.dumps({"range": cell_range, "engine": "openpyxl", "values": rows},
                          ensure_ascii=False, default=str)
    except Exception:
        pass

    if _HAS_PYWIN32:
        try:
            excel = _get_excel()
            wb = excel.Workbooks.Open(
                path, UpdateLinks=0, ReadOnly=True, IgnoreReadOnlyRecommended=True,
                Notify=False, AddToMru=False
            )
            try:
                ws = wb.Worksheets(sheet) if sheet else wb.ActiveSheet
                val = ws.Range(cell_range).Value
                # COM은 2D 튜플 반환 → 리스트로 정규화
                if isinstance(val, tuple):
                    rows = [list(r) if isinstance(r, tuple) else [r] for r in val]
                else:
                    rows = [[val]]
            finally:
                wb.Close(SaveChanges=0)
            if not any(v is None for r in rows for v in r):
                return json.dumps({"range": cell_range, "engine": "com", "values": rows},
                                  ensure_ascii=False, default=str)
        except Exception:
            pass

    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, data_only=True)
        ws = wb[sheet] if sheet else wb.active
        rows = [[c.value for c in row] for row in ws[cell_range]]
        # data_only=True는 Excel이 계산해 캐시한 값만 반환 → openpyxl이 쓴 수식 셀은
        # 캐시가 없어 None이 된다. 그 경우 수식 문자열로 폴백해 의미 있는 값을 돌려준다.
        if any(v is None for r in rows for v in r):
            wbf = openpyxl.load_workbook(path, data_only=False)
            wsf = wbf[sheet] if sheet else wbf.active
            frows = [[c.value for c in row] for row in wsf[cell_range]]
            rows = [
                [v if v is not None else fv for v, fv in zip(r, fr)]
                for r, fr in zip(rows, frows)
            ]
        return json.dumps({"range": cell_range, "engine": "openpyxl", "values": rows},
                          ensure_ascii=False, default=str)
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False)


def excel_active_set_cells(cells: dict, sheet: str = "") -> str:
    """현재 열려있는 활성 Excel 창의 특정 셀들에 값을 실시간으로 설정합니다."""
    if not isinstance(cells, dict) or not cells:
        return json.dumps({"error": "cells는 비어있지 않은 {셀주소: 값} 객체여야 합니다."}, ensure_ascii=False)
    if not _HAS_PYWIN32:
        return _no_com_msg("활성 Excel 제어")
    
    try:
        import win32com.client as win32
        try:
            excel = win32.GetActiveObject("Excel.Application")
        except Exception:
            return json.dumps({"error": "현재 열려있는 Excel 창이 없습니다. 사용자가 Excel을 먼저 열어두어야 합니다."}, ensure_ascii=False)
        
        wb = excel.ActiveWorkbook
        if wb is None:
            return json.dumps({"error": "열려있는 통합 문서가 없습니다."}, ensure_ascii=False)
        ws = wb.Worksheets(sheet) if sheet else wb.ActiveSheet
        for addr, val in cells.items():
            ws.Range(addr).Value = val
        return json.dumps({
            "engine": "com_active", 
            "cells_set": len(cells),
            "message": "활성 Excel 창 실시간 편집 완료 (화면에서 변경사항을 확인하세요)"
        }, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"error": f"활성 문서 제어 오류: {str(e)}"}, ensure_ascii=False)


def excel_active_get_range(cell_range: str, sheet: str = "") -> str:
    """현재 열려있는 활성 Excel 창의 특정 범위 값을 실시간으로 읽습니다."""
    if not _HAS_PYWIN32:
        return _no_com_msg("활성 Excel 제어")
    
    try:
        import win32com.client as win32
        try:
            excel = win32.GetActiveObject("Excel.Application")
        except Exception:
            return json.dumps({"error": "현재 열려있는 Excel 창이 없습니다."}, ensure_ascii=False)
        
        wb = excel.ActiveWorkbook
        if wb is None:
            return json.dumps({"error": "열려있는 통합 문서가 없습니다."}, ensure_ascii=False)
        ws = wb.Worksheets(sheet) if sheet else wb.ActiveSheet
        val = ws.Range(cell_range).Value
        if isinstance(val, tuple):
            rows = [list(r) if isinstance(r, tuple) else [r] for r in val]
        else:
            rows = [[val]]
        return json.dumps({"range": cell_range, "engine": "com_active", "values": rows}, ensure_ascii=False, default=str)
    except Exception as e:
        return json.dumps({"error": f"활성 문서 제어 오류: {str(e)}"}, ensure_ascii=False)



# ── PowerPoint 편집 ───────────────────────────────────────────

def ppt_replace_text(path: str, find: str, replace: str) -> str:
    """PowerPoint(.pptx) 모든 슬라이드에서 텍스트를 찾아 바꾸고 저장합니다(python-pptx).
    편집 전 자동 백업."""
    path = _abspath(path)
    err = _validate_ooxml(path)
    if err:
        return json.dumps({"error": err}, ensure_ascii=False)
    backup = _backup(path)
    try:
        from pptx import Presentation
        prs = Presentation(path)
        n = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    # 1) run 단위 치환(서식 보존)
                    for run in para.runs:
                        if find in run.text:
                            n += run.text.count(find)
                            run.text = run.text.replace(find, replace)
                    # 2) run 경계에 걸친 경우: 단락 전체에서 치환(서식 단순화)
                    joined = "".join(r.text for r in para.runs)
                    if find in joined and para.runs:
                        n += joined.count(find)
                        para.runs[0].text = joined.replace(find, replace)
                        for r in para.runs[1:]:
                            r.text = ""
        prs.save(path)
        return json.dumps({"path": path, "engine": "python-pptx", "backup": backup,
                           "occurrences": n, "message": "PPT 찾아바꾸기 완료"}, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"error": str(e), "backup": backup}, ensure_ascii=False)


def ppt_add_slide(path: str, title: str = "", body: str = "", layout: int = 1) -> str:
    """PowerPoint에 슬라이드를 추가합니다(제목 + 줄바꿈 구분 본문). 파일이 없으면 새로 생성.
    layout: 0=제목, 1=제목+내용(기본), 5=제목만, 6=빈 화면."""
    path = _abspath(path)
    if os.path.exists(path):
        err = _validate_ooxml(path)
        if err:
            return json.dumps({"error": err}, ensure_ascii=False)
    backup = _backup(path)
    try:
        from pptx import Presentation
        prs = Presentation(path) if os.path.exists(path) else Presentation()
        layout = max(0, min(layout, len(prs.slide_layouts) - 1))
        slide = prs.slides.add_slide(prs.slide_layouts[layout])
        if title and slide.shapes.title is not None:
            slide.shapes.title.text = title
        if body:
            # 본문 플레이스홀더 탐색
            ph = None
            for shape in slide.placeholders:
                if shape.placeholder_format.idx != 0 and shape.has_text_frame:
                    ph = shape
                    break
            if ph is not None:
                lines = body.split("\n")
                tf = ph.text_frame
                tf.text = lines[0]
                for line in lines[1:]:
                    p = tf.add_paragraph()
                    p.text = line
        Path(path).parent.mkdir(parents=True, exist_ok=True)
        prs.save(path)
        return json.dumps({"path": path, "engine": "python-pptx", "backup": backup,
                           "slides": len(prs.slides),
                           "message": "슬라이드 추가 완료"}, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"error": str(e), "backup": backup}, ensure_ascii=False)


def ppt_export_pdf(path: str, pdf_path: str = "") -> str:
    """PowerPoint를 PDF로 내보냅니다. (설치된 PowerPoint 필요)"""
    path = _abspath(path)
    err = _validate_ooxml(path)
    if err:
        return json.dumps({"error": err}, ensure_ascii=False)
    if not pdf_path:
        pdf_path = str(Path(path).with_suffix(".pdf"))
    pdf_path = _abspath(pdf_path)
    if _HAS_PYWIN32:
        try:
            ppt = _get_ppt()
            pres = ppt.Presentations.Open(path, ReadOnly=True, WithWindow=False)
            try:
                pres.SaveAs(pdf_path, 32)  # ppSaveAsPDF=32
            finally:
                pres.Close()
            return json.dumps({"path": pdf_path, "engine": "com",
                               "message": "PPT→PDF 내보내기 완료"}, ensure_ascii=False)
        except Exception as e:
            com_err = str(e)
    else:
        com_err = "pywin32/COM 미사용"
    # COM 불가/실패 → LibreOffice 헤드리스 폴백
    out = _libre_to_pdf(path, pdf_path)
    if out:
        return json.dumps({"path": out, "engine": "libreoffice",
                           "message": "PPT→PDF 내보내기 완료(LibreOffice 폴백)"}, ensure_ascii=False)
    return json.dumps({"error": f"PPT→PDF 내보내기 실패: {com_err}. LibreOffice도 사용 불가.",
                       "engine": "none"}, ensure_ascii=False)


def office_close() -> str:
    """열려있는 Word/Excel/PowerPoint COM 세션을 종료합니다(좀비 프로세스 정리)."""
    if not _HAS_PYWIN32:
        return json.dumps({"message": "COM 세션 없음"}, ensure_ascii=False)
    _quit_apps()
    return json.dumps({"message": "Office COM 세션 종료 완료"}, ensure_ascii=False)


# ── MANIFEST ──────────────────────────────────────────────────

MANIFEST = [
    {
        "name": "word_edit_text",
        "label": "Word 찾아바꾸기",
        "schema": {
            "type": "function",
            "function": {
                "name": "word_edit_text",
                "description": (
                    "기존 Word(.docx) 문서에서 텍스트를 찾아 바꾸고 저장합니다(서식 보존). "
                    "설치된 Word(COM) 우선, 불가 시 python-docx 폴백. 편집 전 자동 백업(.bak). "
                    "기존 문서를 실제로 편집할 때 사용하세요(append_word/write_word는 추가/새작성용)."
                ),
                "parameters": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": ".docx 경로"},
                        "find": {"type": "string", "description": "찾을 텍스트"},
                        "replace": {"type": "string", "description": "바꿀 텍스트"},
                        "match_case": {"type": "boolean", "description": "대소문자 구분"},
                        "whole_word": {"type": "boolean", "description": "단어 단위 일치"},
                    },
                    "required": ["path", "find", "replace"],
                },
            },
        },
        "handler": lambda a: word_edit_text(a["path"], a["find"], a["replace"],
                                            a.get("match_case", False), a.get("whole_word", False)),
    },
    {
        "name": "word_insert_text",
        "label": "Word 텍스트 삽입",
        "schema": {
            "type": "function",
            "function": {
                "name": "word_insert_text",
                "description": (
                    "기존 Word 문서에 텍스트를 삽입하고 저장합니다. after_anchor가 있으면 그 텍스트 뒤에, "
                    "없으면 문서 끝에 삽입. RAG/URL 조사 내용을 기존 보고서에 채워 넣을 때 사용. "
                    "삽입할 때 출처(URL/노트 경로)를 함께 적으면 좋습니다. COM 우선, python-docx 폴백(끝에만)."
                ),
                "parameters": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": ".docx 경로"},
                        "text": {"type": "string", "description": "삽입할 텍스트"},
                        "after_anchor": {"type": "string", "description": "이 텍스트 뒤에 삽입(생략 시 문서 끝)"},
                    },
                    "required": ["path", "text"],
                },
            },
        },
        "handler": lambda a: word_insert_text(a["path"], a["text"], a.get("after_anchor", "")),
    },
    {
        "name": "word_export_pdf",
        "label": "Word→PDF 내보내기",
        "schema": {
            "type": "function",
            "function": {
                "name": "word_export_pdf",
                "description": "Word(.docx)를 PDF로 내보냅니다. 설치된 Word(COM) 필요.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": ".docx 경로"},
                        "pdf_path": {"type": "string", "description": "출력 PDF 경로(생략 시 같은 이름)"},
                    },
                    "required": ["path"],
                },
            },
        },
        "handler": lambda a: word_export_pdf(a["path"], a.get("pdf_path", "")),
    },
    {
        "name": "word_accept_all_changes",
        "label": "Word 수정추적 수락",
        "schema": {
            "type": "function",
            "function": {
                "name": "word_accept_all_changes",
                "description": "Word 문서의 모든 수정추적(Track Changes)을 수락하고 저장합니다. 설치된 Word(COM) 필요.",
                "parameters": {
                    "type": "object",
                    "properties": {"path": {"type": "string", "description": ".docx 경로"}},
                    "required": ["path"],
                },
            },
        },
        "handler": lambda a: word_accept_all_changes(a["path"]),
    },
    {
        "name": "word_add_comment",
        "label": "Word 메모 작성",
        "schema": {
            "type": "function",
            "function": {
                "name": "word_add_comment",
                "description": "Word 문서에서 특정 텍스트(anchor_text)에 검토 메모를 추가합니다. 설치된 Word(COM) 필요.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": ".docx 경로"},
                        "anchor_text": {"type": "string", "description": "메모를 달 기준 텍스트"},
                        "comment": {"type": "string", "description": "메모 내용"},
                    },
                    "required": ["path", "anchor_text", "comment"],
                },
            },
        },
        "handler": lambda a: word_add_comment(a["path"], a["anchor_text"], a["comment"]),
    },
    {
        "name": "excel_set_cells",
        "label": "Excel 셀 편집",
        "schema": {
            "type": "function",
            "function": {
                "name": "excel_set_cells",
                "description": (
                    "기존 Excel(.xlsx) 워크북의 특정 셀들에 값/수식을 설정하고 저장합니다. "
                    "cells 예: {\"B2\": \"=A1+A2\", \"C3\": \"합계\"}. "
                    "설치된 Excel(COM) 우선(수식 재계산·서식 보존), 불가 시 openpyxl 폴백. 편집 전 자동 백업."
                ),
                "parameters": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": ".xlsx 경로"},
                        "cells": {"type": "object", "description": "{셀주소: 값/수식} 객체"},
                        "sheet": {"type": "string", "description": "시트 이름(생략 시 활성 시트)"},
                    },
                    "required": ["path", "cells"],
                },
            },
        },
        "handler": lambda a: excel_set_cells(a["path"], a["cells"], a.get("sheet", "")),
    },
    {
        "name": "excel_get_range",
        "label": "Excel 범위 읽기",
        "schema": {
            "type": "function",
            "function": {
                "name": "excel_get_range",
                "description": "Excel 워크북의 특정 범위(예 'A1:C10') 값을 읽습니다. COM 우선, openpyxl 폴백.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": ".xlsx 경로"},
                        "cell_range": {"type": "string", "description": "예: 'A1:C10'"},
                        "sheet": {"type": "string", "description": "시트 이름(생략 시 활성 시트)"},
                    },
                    "required": ["path", "cell_range"],
                },
            },
        },
        "handler": lambda a: excel_get_range(a["path"], a["cell_range"], a.get("sheet", "")),
    },
    {
        "name": "excel_active_set_cells",
        "label": "실시간 Excel 셀 편집",
        "schema": {
            "type": "function",
            "function": {
                "name": "excel_active_set_cells",
                "description": (
                    "현재 사용자가 열어둔 활성 Excel 창의 특정 셀들에 값을 실시간으로 설정합니다. "
                    "사용자가 눈으로 직접 엑셀을 보면서 호흡을 맞추며 작업할 때 사용하세요. "
                    "파일 경로 없이 현재 화면의 엑셀에 바로 입력합니다. 저장은 하지 않습니다."
                ),
                "parameters": {
                    "type": "object",
                    "properties": {
                        "cells": {"type": "object", "description": "셀 주소와 값 딕셔너리. 예: {\"A1\": 100, \"B1\": \"=A1*2\"}"},
                        "sheet": {"type": "string", "description": "시트명 (생략 시 현재 활성 시트)"},
                    },
                    "required": ["cells"],
                },
            },
        },
        "handler": lambda a: excel_active_set_cells(a["cells"], a.get("sheet", "")),
    },
    {
        "name": "excel_active_get_range",
        "label": "실시간 Excel 셀 읽기",
        "schema": {
            "type": "function",
            "function": {
                "name": "excel_active_get_range",
                "description": (
                    "현재 사용자가 열어둔 활성 Excel 창의 특정 범위 값을 실시간으로 읽습니다. "
                    "파일 경로 없이 현재 화면에 켜져 있는 엑셀 문서를 기준으로 합니다."
                ),
                "parameters": {
                    "type": "object",
                    "properties": {
                        "cell_range": {"type": "string", "description": "범위 주소. 예: 'A1:C10'"},
                        "sheet": {"type": "string", "description": "시트명 (생략 시 현재 활성 시트)"},
                    },
                    "required": ["cell_range"],
                },
            },
        },
        "handler": lambda a: excel_active_get_range(a["cell_range"], a.get("sheet", "")),
    },
    {
        "name": "ppt_replace_text",
        "label": "PPT 찾아바꾸기",
        "schema": {
            "type": "function",
            "function": {
                "name": "ppt_replace_text",
                "description": "PowerPoint(.pptx) 모든 슬라이드에서 텍스트를 찾아 바꾸고 저장합니다. 편집 전 자동 백업.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": ".pptx 경로"},
                        "find": {"type": "string"},
                        "replace": {"type": "string"},
                    },
                    "required": ["path", "find", "replace"],
                },
            },
        },
        "handler": lambda a: ppt_replace_text(a["path"], a["find"], a["replace"]),
    },
    {
        "name": "ppt_add_slide",
        "label": "PPT 슬라이드 추가",
        "schema": {
            "type": "function",
            "function": {
                "name": "ppt_add_slide",
                "description": ("PowerPoint에 슬라이드를 추가합니다(제목 + 줄바꿈 구분 본문). "
                                "파일이 없으면 새로 생성. layout: 1=제목+내용(기본)."),
                "parameters": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": ".pptx 경로"},
                        "title": {"type": "string", "description": "슬라이드 제목"},
                        "body": {"type": "string", "description": "본문(줄바꿈으로 단락 구분)"},
                        "layout": {"type": "integer", "description": "0=제목,1=제목+내용,5=제목만,6=빈화면"},
                    },
                    "required": ["path"],
                },
            },
        },
        "handler": lambda a: ppt_add_slide(a["path"], a.get("title", ""), a.get("body", ""), a.get("layout", 1)),
    },
    {
        "name": "ppt_export_pdf",
        "label": "PPT→PDF 내보내기",
        "schema": {
            "type": "function",
            "function": {
                "name": "ppt_export_pdf",
                "description": "PowerPoint(.pptx)를 PDF로 내보냅니다. 설치된 PowerPoint(COM) 필요.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": ".pptx 경로"},
                        "pdf_path": {"type": "string", "description": "출력 PDF 경로(생략 시 같은 이름)"},
                    },
                    "required": ["path"],
                },
            },
        },
        "handler": lambda a: ppt_export_pdf(a["path"], a.get("pdf_path", "")),
    },
    {
        "name": "office_close",
        "label": "Office 세션 종료",
        "schema": {
            "type": "function",
            "function": {
                "name": "office_close",
                "description": "열려있는 Word/Excel COM 세션을 종료합니다(좀비 프로세스 정리). 편집 작업을 모두 마친 뒤 호출하세요.",
                "parameters": {"type": "object", "properties": {}},
            },
        },
        "handler": lambda a: office_close(),
    },
]

# 모든 Office 핸들러를 전용 COM(STA) 스레드로 위임 — 단일 스레드 직렬화로 COM 안정성 확보
# + OFFICE_COM_TIMEOUT 워치독으로 무한 행 방지(타임아웃 시 인스턴스 정리 후 구조화 오류 반환)
for _tool in MANIFEST:
    _tool["handler"] = _on_com_thread(_tool["handler"], _tool["name"])
